# code/converters/markdown_converter.py
"""
Конвертеры различных форматов в Markdown
"""
import os
import asyncio
import logging
import tempfile
import subprocess
import shutil
from typing import Tuple, Optional, Any, List
from io import BytesIO
import pandas as pd
from PIL import Image
from converters.base_converter import FileConverter, ProgressCallback, CancelCheck

logger = logging.getLogger(__name__)


class MarkdownConverterManager:
    """Менеджер конвертеров в Markdown"""
    
    def __init__(self, ai_interface=None):
        self.ai_interface = ai_interface
        self.converters = [
            PDFConverter(ai_interface),
            ImageConverter(ai_interface),
            ExcelConverter(),
            CSVConverter(),
            DocumentConverter(),
            HTMLConverter(),
            PowerPointConverter()
        ]
    
    async def convert_to_markdown(
        self, 
        user_id: str,
        file_bytes: bytes, 
        filename: str,
        encoding: Optional[str] = None,
        progress_callback: Optional[ProgressCallback] = None,
        cancel_check: Optional[CancelCheck] = None
    ) -> Tuple[bool, str, str]:
        """
        Конвертирует файл в Markdown если возможно
        Returns: (success, new_filename, markdown_content)
        """
        for converter in self.converters:
            if converter.supports(filename):
                kwargs = {}
                if encoding:
                    kwargs['encoding'] = encoding
                if progress_callback:
                    kwargs['progress_callback'] = progress_callback
                if cancel_check:
                    kwargs['cancel_check'] = cancel_check
                
                return await converter.convert(user_id, file_bytes, filename, **kwargs)
        
        return False, filename, ""
    
    def can_convert(self, filename: str) -> bool:
        """Проверяет, можно ли конвертировать файл"""
        return any(conv.supports(filename) for conv in self.converters)


class HTMLConverter(FileConverter):
    """Конвертер HTML в Markdown используя markdownify"""
    
    SUPPORTED_EXTENSIONS = {'.html', '.htm', '.xhtml'}
    
    def supports(self, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.SUPPORTED_EXTENSIONS
    
    async def convert(self, user_id: str, file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Конвертирует HTML в Markdown"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        try:
            # Пробуем импортировать markdownify
            try:
                from markdownify import markdownify
            except ImportError:
                logger.warning("markdownify не установлен, пробуем альтернативный метод")
                return await self._convert_via_pandoc(file_bytes, filename)
            
            # Определяем кодировку HTML
            encoding = kwargs.get('encoding', 'utf-8')
            try:
                html_content = file_bytes.decode(encoding)
            except UnicodeDecodeError:
                # Пробуем другие кодировки
                for enc in ['utf-8', 'cp1251', 'latin-1', 'cp866']:
                    try:
                        html_content = file_bytes.decode(enc)
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    html_content = file_bytes.decode('utf-8', errors='replace')
            
            # Конвертируем HTML в Markdown
            markdown_text = markdownify(
                html_content,
                heading_style="ATX",  # Используем # для заголовков
                bullets="-",  # Используем - для списков
                code_language="python",  # Язык по умолчанию для блоков кода
                strip=['script', 'style', 'meta', 'link'],  # Удаляем эти теги
                wrap=True,
                wrap_width=80,
                convert=['a', 'b', 'blockquote', 'br', 'code', 'div', 'em', 
                        'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'hr', 'i', 
                        'img', 'li', 'ol', 'p', 'pre', 'strong', 'u', 'ul',
                        'table', 'thead', 'tbody', 'tr', 'th', 'td']
            )
            
            # Очищаем результат
            markdown_text = self._clean_markdown(markdown_text)
            
            if markdown_text and markdown_text.strip():
                logger.info(f"HTML успешно конвертирован через markdownify: {filename}")
                return True, new_name, markdown_text
            else:
                logger.warning(f"markdownify вернул пустой результат для {filename}")
                return False, new_name, ""
                
        except Exception as e:
            logger.error(f"Ошибка конвертации HTML через markdownify: {e}")
            # Пробуем через pandoc как запасной вариант
            return await self._convert_via_pandoc(file_bytes, filename)
    
    async def _convert_via_pandoc(self, file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
        """Альтернативная конвертация через pandoc"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        if not shutil.which("pandoc"):
            logger.error("pandoc не установлен")
            return False, new_name, ""
        
        try:
            with tempfile.TemporaryDirectory() as td:
                in_path = os.path.join(td, "input.html")
                out_path = os.path.join(td, "output.md")
                
                with open(in_path, "wb") as f:
                    f.write(file_bytes)
                
                # Используем pandoc с параметрами для хорошего markdown
                subprocess.run(
                    ["pandoc", "-s", in_path, "--wrap=none", 
                     "--reference-links", "-t", "markdown", "-o", out_path],
                    check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
                )
                
                with open(out_path, "r", encoding="utf-8", errors="replace") as f:
                    markdown_text = f.read()
            
            markdown_text = self._clean_markdown(markdown_text)
            logger.info(f"HTML успешно конвертирован через pandoc: {filename}")
            return True, new_name, markdown_text
            
        except Exception as e:
            logger.error(f"Ошибка pandoc конвертации HTML: {e}")
            return False, new_name, ""
    
    def _clean_markdown(self, text: str) -> str:
        """Очищает и форматирует markdown текст"""
        import re
        
        # Удаляем множественные пустые строки
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Удаляем пробелы в начале и конце строк
        lines = text.split('\n')
        lines = [line.rstrip() for line in lines]
        text = '\n'.join(lines)
        
        # Удаляем HTML комментарии
        text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
        
        # Удаляем оставшиеся HTML теги
        text = re.sub(r'<[^>]+>', '', text)
        
        return text.strip()


class PowerPointConverter(FileConverter):
    """Конвертер PowerPoint в Markdown используя pptx2md"""
    
    SUPPORTED_EXTENSIONS = {'.pptx', '.ppt'}
    
    def supports(self, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.SUPPORTED_EXTENSIONS
    
    async def convert(self, user_id: str, file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Конвертирует PowerPoint в Markdown"""
        base, ext = os.path.splitext(filename)
        new_name = f"{base}.txt"
    
        # Для .ppt файлов пробуем только pandoc
        if ext.lower() == '.ppt':
            return await self._convert_via_pandoc(file_bytes, filename)
    
        # Сразу используем python-pptx, который работает стабильно
        return await self._convert_via_python_pptx(file_bytes, filename)
    
    async def _convert_via_python_pptx(self, file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
        """Альтернативная конвертация через python-pptx"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx не установлен")
            return await self._convert_via_pandoc(file_bytes, filename)
        
        try:
            prs = Presentation(BytesIO(file_bytes))
            markdown_lines = [f"# {filename}\n"]
            
            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_lines.append(f"\n## Слайд {slide_num}\n")
                
                # Извлекаем текст из всех shape
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            # Определяем тип shape для форматирования
                            if shape.name.startswith("Title"):
                                markdown_lines.append(f"### {text}\n")
                            elif shape.name.startswith("Subtitle"):
                                markdown_lines.append(f"**{text}**\n")
                            else:
                                markdown_lines.append(f"{text}\n")
                    
                    # Обрабатываем таблицы
                    if shape.has_table:
                        markdown_lines.append("\n")
                        table = shape.table
                        
                        # Заголовки таблицы
                        headers = []
                        for cell in table.rows[0].cells:
                            headers.append(cell.text.strip())
                        
                        if headers:
                            markdown_lines.append("| " + " | ".join(headers) + " |")
                            markdown_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                            
                            # Строки таблицы
                            for row in table.rows[1:]:
                                cells = []
                                for cell in row.cells:
                                    cells.append(cell.text.strip())
                                markdown_lines.append("| " + " | ".join(cells) + " |")
                        
                        markdown_lines.append("\n")
                
                # Добавляем разделитель между слайдами
                if slide_num < len(prs.slides):
                    markdown_lines.append("\n---\n")
            
            markdown_text = "\n".join(markdown_lines)
            
            if markdown_text.strip():
                logger.info(f"PowerPoint успешно конвертирован через python-pptx: {filename}")
                return True, new_name, markdown_text
            else:
                return False, new_name, ""
                
        except Exception as e:
            logger.error(f"Ошибка конвертации PowerPoint через python-pptx: {e}")
            return await self._convert_via_pandoc(file_bytes, filename)
    
    async def _convert_via_pandoc(self, file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
        """Запасная конвертация через pandoc"""
        base, ext = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        if not shutil.which("pandoc"):
            logger.error("pandoc не установлен")
            return False, new_name, ""
        
        try:
            with tempfile.TemporaryDirectory() as td:
                in_path = os.path.join(td, f"input{ext}")
                out_path = os.path.join(td, "output.md")
                
                with open(in_path, "wb") as f:
                    f.write(file_bytes)
                
                subprocess.run(
                    ["pandoc", "-s", in_path, "--wrap=none", 
                     "-t", "markdown", "-o", out_path],
                    check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
                )
                
                with open(out_path, "r", encoding="utf-8", errors="replace") as f:
                    markdown_text = f.read()
            
            logger.info(f"PowerPoint успешно конвертирован через pandoc: {filename}")
            return True, new_name, markdown_text
            
        except Exception as e:
            logger.error(f"Ошибка pandoc конвертации PowerPoint: {e}")
            return False, new_name, ""


class ExcelConverter(FileConverter):
    """Конвертер Excel в Markdown используя excel-to-markdown"""
    
    def supports(self, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in {'.xls', '.xlsx'}
    
    async def convert(self, user_id: str, file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Конвертирует Excel в Markdown"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        # Сначала пробуем excel-to-markdown если установлен
        result = await self._convert_via_excel_to_markdown(file_bytes, filename)
        if result[0]:
            return result
        
        # Если не получилось, используем pandas
        return await self._convert_via_pandas(file_bytes, filename)
    
    async def _convert_via_excel_to_markdown(self, file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
        """Конвертация через excel-to-markdown"""
        base, ext = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        # Проверяем наличие excel-to-markdown
        if not shutil.which("excel-to-markdown"):
            logger.debug("excel-to-markdown не установлен")
            return False, new_name, ""
        
        try:
            with tempfile.TemporaryDirectory() as td:
                in_path = os.path.join(td, f"input{ext}")
                out_path = os.path.join(td, "output.md")
                
                with open(in_path, "wb") as f:
                    f.write(file_bytes)
                
                # Используем excel-to-markdown
                subprocess.run(
                    ["excel-to-markdown", in_path, "-o", out_path],
                    check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
                )
                
                with open(out_path, "r", encoding="utf-8", errors="replace") as f:
                    markdown_text = f.read()
            
            if markdown_text.strip():
                logger.info(f"Excel успешно конвертирован через excel-to-markdown: {filename}")
                return True, new_name, markdown_text
            else:
                return False, new_name, ""
                
        except Exception as e:
            logger.debug(f"Ошибка excel-to-markdown: {e}")
            return False, new_name, ""
    
    async def _convert_via_pandas(self, file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
        """Конвертация через pandas (fallback)"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        try:
            bio = BytesIO(file_bytes)
            xls = pd.ExcelFile(bio)
            out_parts = [f"# {filename}\n"]
            
            for sheet in xls.sheet_names:
                df = xls.parse(sheet, dtype=object)
                df = df.fillna("")
                
                out_parts.append(f"\n## {sheet}\n")
                
                # Статистика по листу
                out_parts.append(f"*Строк: {len(df)}, Столбцов: {len(df.columns)}*\n\n")
                
                # Конвертируем в markdown таблицу
                rows = [list(df.columns)] + df.values.tolist()
                out_parts.append(self._rows_to_markdown(rows))
                out_parts.append("\n")
            
            md = "\n".join(out_parts).strip() + "\n"
            logger.info(f"Excel успешно конвертирован через pandas: {filename}")
            return True, new_name, md
            
        except Exception as e:
            logger.error(f"Ошибка конвертации Excel через pandas: {e}")
            return False, new_name, ""
    
    def _rows_to_markdown(self, rows: List[List[Any]]) -> str:
        """Преобразует список строк в markdown-таблицу"""
        if not rows:
            return ""
        
        str_rows = [[("" if c is None else str(c)) for c in r] for r in rows]
        header = str_rows[0]
        body = str_rows[1:] if len(str_rows) > 1 else []
        
        widths = [len(h) for h in header]
        for r in body:
            for j, cell in enumerate(r):
                if j < len(widths):
                    widths[j] = max(widths[j], len(cell))
                else:
                    widths.append(len(cell))
        
        def fmt_row(r):
            return "| " + " | ".join((r[i] if i < len(r) else "").ljust(widths[i]) for i in range(len(widths))) + " |"
        
        parts = [fmt_row(header)]
        parts.append("| " + " | ".join("-" * widths[i] for i in range(len(widths))) + " |")
        for r in body:
            parts.append(fmt_row(r))
        
        return "\n".join(parts)


class DocumentConverter(FileConverter):
    """Конвертер документов Word/RTF через pandoc с улучшенными параметрами"""
    
    SUPPORTED_EXTENSIONS = {'.docx', '.doc', '.rtf', '.odt'}
    
    def supports(self, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.SUPPORTED_EXTENSIONS and shutil.which("pandoc") is not None
    
    async def convert(self, user_id: str, file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Конвертирует документ в Markdown через pandoc"""
        base, ext = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        try:
            with tempfile.TemporaryDirectory() as td:
                in_path = os.path.join(td, f"input{ext}")
                out_path = os.path.join(td, "output.md")
                
                with open(in_path, "wb") as f:
                    f.write(file_bytes)
                
                # Используем расширенные параметры pandoc для лучшего markdown
                pandoc_args = [
                    "pandoc",
                    "-s",  # standalone
                    in_path,
                    "--wrap=auto",  # автоматический перенос строк
                    "--columns=128",  # ширина колонки 128 символов
                    "--reference-links",  # ссылки в конце документа
                    "-t", "markdown",  # выходной формат
                    "--extract-media", td,  # извлечь медиа в папку
                    "-o", out_path
                ]
                
                # Для docx можно добавить дополнительные опции
                if ext.lower() == '.docx':
                    pandoc_args.insert(2, "--track-changes=accept")  # принять все изменения
                
                result = subprocess.run(
                    pandoc_args,
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                
                if result.returncode != 0:
                    logger.warning(f"pandoc вернул код {result.returncode}: {result.stderr}")
                
                with open(out_path, "r", encoding="utf-8", errors="replace") as f:
                    md = f.read()
            
            # Постобработка markdown
            md = self._postprocess_markdown(md)
            
            if md.strip():
                logger.info(f"Документ успешно конвертирован через pandoc: {filename}")
                return True, new_name, md
            else:
                logger.warning(f"pandoc вернул пустой результат для {filename}")
                return False, new_name, ""
            
        except subprocess.TimeoutExpired:
            logger.error(f"Таймаут при конвертации документа: {filename}")
            return False, new_name, ""
        except Exception as e:
            logger.error(f"Ошибка pandoc конвертации: {e}")
            return False, new_name, ""
    
    def _postprocess_markdown(self, text: str) -> str:
        """Постобработка markdown после pandoc"""
        import re
    
        # Удаляем разметку {.mark} и подобную
        text = re.sub(r'\[([^\]]+)\]\{[^}]+\}', r'\1', text)
        text = re.sub(r'\{\.[\w\-]+\}', '', text)
    
        # Удаляем escape-символы перед специальными символами markdown
        text = re.sub(r'\\([#*_\[\]()>"|])', r'\1', text)
    
        # Удаляем двойные обратные слеши
        text = re.sub(r'\\\\', '', text)
    
        # Удаляем лишние пустые строки
        text = re.sub(r'\n{3,}', '\n\n', text)
    
        # Удаляем пробелы в конце строк
        lines = text.split('\n')
        lines = [line.rstrip() for line in lines]
        text = '\n'.join(lines)
    
        # Исправляем форматирование таблиц
        text = self._fix_tables(text)
    
        # Удаляем мета-теги в начале документа если есть
        text = re.sub(r'^---\n.*?\n---\n', '', text, flags=re.DOTALL)
    
        # Удаляем span-теги и их содержимое в квадратных скобках
        text = re.sub(r'\[([^\]]+)\]\([^)]*\)', r'\1', text)
    
        # Убираем лишние квадратные скобки вокруг текста
        text = re.sub(r'^\[(.+)\]$', r'\1', text, flags=re.MULTILINE)
    
        return text.strip()
    
    def _fix_tables(self, text: str) -> str:
        """Исправляет форматирование таблиц"""
        import re
    
        lines = text.split('\n')
        fixed_lines = []
        in_table = False
    
        for line in lines:
            if '|' in line:
                if not in_table:
                    in_table = True
                # Нормализуем разделители в таблице
                if re.match(r'^[\s\|:\-]+$', line):
                    # Это строка-разделитель
                    parts = line.split('|')
                    parts = ['---' if p.strip() else '' for p in parts]
                    line = '|'.join(parts)
                fixed_lines.append(line)
            else:
                if in_table:
                    in_table = False
                fixed_lines.append(line)
    
        return '\n'.join(fixed_lines)


class CSVConverter(FileConverter):
    """Конвертер CSV в Markdown"""
    
    def supports(self, filename: str) -> bool:
        return filename.lower().endswith('.csv')
    
    async def convert(self, user_id: str, file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Конвертирует CSV в Markdown"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        encoding = kwargs.get('encoding', 'utf-8')
        
        try:
            bio = BytesIO(file_bytes)
            df = pd.read_csv(bio, encoding=encoding, dtype=object)
        except Exception:
            # Пробуем определить кодировку
            import chardet
            result = chardet.detect(file_bytes)
            detected = result.get('encoding', 'utf-8')
            try:
                bio = BytesIO(file_bytes)
                df = pd.read_csv(bio, encoding=detected, dtype=object)
            except Exception as e:
                logger.error(f"Ошибка конвертации CSV: {e}")
                return False, new_name, ""
        
        df = df.fillna("")
        rows = [list(df.columns)] + df.values.tolist()
        
        # Используем тот же метод для форматирования
        converter = ExcelConverter()
        md = f"# {filename}\n\n"
        md += f"*Строк: {len(df)}, Столбцов: {len(df.columns)}*\n\n"
        md += converter._rows_to_markdown(rows) + "\n"
        
        return True, new_name, md

class PDFConverter(FileConverter):
    """Конвертер PDF в Markdown"""
    
    def __init__(self, ai_interface=None):
        self.ai_interface = ai_interface
        self._markitdown = None
    
    def supports(self, filename: str) -> bool:
        return filename.lower().endswith('.pdf')
    
    def _get_markitdown(self):
        """Lazy инициализация MarkItDown"""
        if self._markitdown is None:
            try:
                from markitdown import MarkItDown
                self._markitdown = MarkItDown()
                logger.info("MarkItDown инициализирован")
            except ImportError:
                logger.warning("MarkItDown не установлен")
                return None
        return self._markitdown
    
    async def convert(self, user_id: str, file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Конвертирует PDF в Markdown"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        progress_callback = kwargs.get('progress_callback')
        cancel_check = kwargs.get('cancel_check')
        
        # Сначала пробуем MarkItDown
        ok, _, md_text = await self._convert_via_markitdown(file_bytes, filename)
        if ok:
            return ok, new_name, md_text
        
        # Затем пробуем OCR (если PDF содержит сканированные изображения)
        ok, _, md_text = await self._convert_via_ocr(file_bytes, filename, progress_callback, cancel_check)
        if ok:
            return ok, new_name, md_text
        
        # Если не получилось и есть AI - используем распознавание через изображения
        if self.ai_interface:
            ok, _, md_text = await self._convert_via_images(
                user_id,
                file_bytes, filename, progress_callback, cancel_check
            )
        
        return ok, new_name, md_text
    
    async def _convert_via_markitdown(self, file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
        """Конвертирует через MarkItDown"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        markitdown = self._get_markitdown()
        if not markitdown:
            return False, new_name, ""
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
                tmp_file.write(file_bytes)
                tmp_path = tmp_file.name
            
            try:
                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(None, markitdown.convert, tmp_path)
                
                if result and hasattr(result, 'text_content'):
                    markdown_text = result.text_content
                    if markdown_text and markdown_text.strip():
                        logger.info(f"PDF успешно конвертирован через MarkItDown: {filename}")
                        return True, new_name, markdown_text
            finally:
                try:
                    os.unlink(tmp_path)
                except:
                    pass
                    
        except Exception as e:
            logger.warning(f"MarkItDown не смог обработать PDF: {e}")
        
        return False, new_name, ""
    
    async def _convert_via_ocr(self, file_bytes: bytes, filename: str, 
                               progress_callback: Optional[ProgressCallback] = None,
                               cancel_check: Optional[CancelCheck] = None) -> Tuple[bool, str, str]:
        """OCR для сканированных PDF с использованием pytesseract"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        try:
            import pytesseract
            from PIL import Image
            from pdf2image import convert_from_bytes
        except ImportError as e:
            logger.debug(f"OCR библиотеки не установлены: {e}")
            return False, new_name, ""
        
        try:
            if progress_callback:
                await progress_callback("🔍 Пробую OCR распознавание PDF...")
            
            logger.info(f"Attempting OCR for potential scanned PDF: {filename}")
            
            # Конвертируем PDF в изображения
            images = convert_from_bytes(file_bytes, dpi=150)
            total_pages = len(images)
            
            if progress_callback:
                await progress_callback(
                    f"📊 Обнаружено страниц: {total_pages}\n"
                    f"⏳ Начинаю OCR распознавание...\n"
                    f"💡 Отправьте /cancel_file_job для отмены"
                )
            
            text_parts = []
            
            for i, image in enumerate(images, 1):
                if cancel_check and await cancel_check():
                    logger.info(f"OCR отменено на странице {i}")
                    if progress_callback:
                        await progress_callback(
                            f"⚠️ OCR распознавание отменено\n"
                            f"✅ Обработано страниц: {i-1}/{total_pages}"
                        )
                    break
                
                if progress_callback:
                    await progress_callback(
                        f"📄 OCR страница {i}/{total_pages}\n"
                        f"{"▓" * i}{"░" * (total_pages - i)} {i}/{total_pages}\n"
                        f"💡 /cancel_file_job для остановки"
                    )
                
                try:
                    # Применяем OCR с поддержкой русского и английского языков
                    text = pytesseract.image_to_string(image, lang='rus+eng')
                    if text and text.strip():
                        text_parts.append(f"## Страница {i}\n\n{text.strip()}")
                        logger.info(f"OCR страница {i} успешно распознана")
                except Exception as e:
                    logger.warning(f"OCR failed for page {i}: {e}")
                    # Пробуем только английский язык
                    try:
                        text = pytesseract.image_to_string(image)
                        if text and text.strip():
                            text_parts.append(f"## Страница {i}\n\n{text.strip()}")
                            logger.info(f"OCR страница {i} распознана (только eng)")
                    except:
                        pass
                
                if i < total_pages:
                    await asyncio.sleep(0.1)  # Небольшая пауза между страницами
            
            if text_parts:
                full_text = f"# Документ: {filename}\n"
                full_text += f"Страниц распознано: {len(text_parts)} из {total_pages}\n"
                full_text += "---\n\n"
                full_text += "\n\n".join(text_parts)
                
                if progress_callback:
                    await progress_callback(
                        f"✅ OCR распознавание завершено!\n"
                        f"📊 Обработано страниц: {len(text_parts)}/{total_pages}"
                    )
                
                logger.info(f"PDF успешно распознан через OCR: {len(text_parts)}/{total_pages} страниц")
                return True, new_name, full_text
                
        except ImportError:
            logger.debug("OCR инструменты не доступны (pytesseract/pdf2image)")
        except Exception as e:
            logger.warning(f"OCR failed for {filename}: {e}")
            if progress_callback:
                await progress_callback(f"⚠️ Ошибка OCR: {str(e)}")
        
        return False, new_name, ""
    
    async def _convert_via_images(
        self, 
        user_id: str,
        file_bytes: bytes, 
        filename: str,
        progress_callback: Optional[ProgressCallback] = None,
        cancel_check: Optional[CancelCheck] = None
    ) -> Tuple[bool, str, str]:
        """Конвертирует PDF через распознавание изображений страниц с помощью AI"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        if not self.ai_interface:
            return False, new_name, ""
        
        try:
            from pdf2image import convert_from_bytes
        except ImportError:
            logger.error("pdf2image не установлен")
            return False, new_name, ""
        
        try:
            if progress_callback:
                await progress_callback("📄 Анализирую PDF файл...")
            
            images = convert_from_bytes(file_bytes, dpi=150)
            total_pages = len(images)
            
            if progress_callback:
                await progress_callback(
                    f"📊 Обнаружено страниц: {total_pages}\n"
                    f"⏳ Начинаю распознавание...\n"
                    f"💡 Отправьте /cancel_file_job для отмены"
                )
            
            page_texts = []
            
            for i, img in enumerate(images):
                page_num = i + 1
                
                if cancel_check and await cancel_check():
                    logger.info(f"Распознавание отменено на странице {page_num}")
                    if progress_callback:
                        await progress_callback(
                            f"⚠️ Распознавание отменено\n"
                            f"✅ Обработано страниц: {len(page_texts)}/{total_pages}"
                        )
                    break
                
                if progress_callback:
                    await progress_callback(
                        f"📄 Распознаю страницу {page_num}/{total_pages}\n"
                        f"{"▓" * page_num}{"░" * (total_pages - page_num)} {page_num}/{total_pages}\n"
                        f"💡 /cancel_file_job для остановки"
                    )
                
                # Конвертируем изображение в bytes
                buffer = BytesIO()
                # Масштабируем если нужно
                width, height = img.size
                if width > 2000 or height > 2000:
                    if width > height:
                        new_width = 2000
                        new_height = int(height * (2000 / width))
                    else:
                        new_height = 2000
                        new_width = int(width * (2000 / height))
                    img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                
                img.save(buffer, format='PNG')
                img_bytes = buffer.getvalue()
                
                prompt = (
                    f"Распознай весь текст на странице {page_num} из {total_pages}. "
                    f"Верни полный текст включая таблицы, списки и все детали в формате Markdown. "
                    f"Если текста нет, напиши только: 'Текст не обнаружен'"
                )
                
                try:
                    result = await self.ai_interface.process_image(user_id, img_bytes, prompt)
                    if result and result.strip() and "Текст не обнаружен" not in result:
                        page_texts.append(f"## Страница {page_num}\n\n{result.strip()}")
                        logger.info(f"Страница {page_num} распознана")
                except Exception as e:
                    logger.error(f"Ошибка распознавания страницы {page_num}: {e}")
                    if progress_callback:
                        await progress_callback(f"⚠️ Ошибка на странице {page_num}, пропускаю...")
                
                if i < len(images) - 1:
                    await asyncio.sleep(0.3)
            
            if page_texts:
                full_text = f"# Документ: {filename}\n"
                full_text += f"Страниц распознано: {len(page_texts)} из {total_pages}\n"
                full_text += "---\n\n"
                full_text += "\n\n".join(page_texts)
                
                if progress_callback:
                    await progress_callback(
                        f"✅ Распознавание завершено!\n"
                        f"📊 Обработано страниц: {len(page_texts)}/{total_pages}"
                    )
                
                logger.info(f"PDF успешно распознан: {len(page_texts)}/{total_pages} страниц")
                return True, new_name, full_text
                
        except Exception as e:
            logger.error(f"Ошибка распознавания PDF: {e}")
            if progress_callback:
                await progress_callback(f"❌ Ошибка распознавания: {str(e)}")
        
        return False, new_name, ""

class ImageConverter(FileConverter):
    """Конвертер изображений в Markdown через AI OCR"""
    
    IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp'}
    
    def __init__(self, ai_interface=None):
        self.ai_interface = ai_interface
    
    def supports(self, filename: str) -> bool:
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.IMAGE_EXTENSIONS and self.ai_interface is not None
    
    async def convert(self, user_id: str,file_bytes: bytes, filename: str, **kwargs) -> Tuple[bool, str, str]:
        """Распознает текст на изображении через AI"""
        base, _ = os.path.splitext(filename)
        new_name = f"{base}.txt"
        
        if not self.ai_interface:
            return False, new_name, ""
        
        try:
            # Масштабируем изображение если нужно
            processed_bytes = await self._resize_image_if_needed(file_bytes)
            
            logger.info(f"Распознаю изображение {filename}, размер: {len(processed_bytes)} байт")
            
            prompt = (
                "Внимательно распознай ВЕСЬ текст на этом изображении. "
                "Верни полный текст включая таблицы, списки и все детали в формате Markdown. "
                "Если текста нет, напиши только: 'Текст не обнаружен'"
            )
            
            result = await self.ai_interface.process_image(user_id,processed_bytes, prompt)
            
            if result:
                result_text = result.strip()
                if result_text and "Текст не обнаружен" not in result_text and "текст не обнаружен" not in result_text.lower():
                    logger.info(f"Изображение распознано, получено {len(result_text)} символов")
                    return True, new_name, result_text
                else:
                    logger.info("Изображение не содержит текста")
                    return False, new_name, ""
            
        except Exception as e:
            logger.error(f"Ошибка AI распознавания изображения: {e}")
        
        return False, new_name, ""
    
    async def _resize_image_if_needed(self, image_bytes: bytes, max_dimension: int = 2000) -> bytes:
        """Масштабирует изображение если оно больше max_dimension"""
        try:
            img = Image.open(BytesIO(image_bytes))
            width, height = img.size
        
            if width <= max_dimension and height <= max_dimension:
                return image_bytes
        
            if width > height:
                new_width = max_dimension
                new_height = int(height * (max_dimension / width))
            else:
                new_height = max_dimension
                new_width = int(width * (max_dimension / height))
        
            logger.info(f"Масштабирую изображение с {width}x{height} до {new_width}x{new_height}")
        
            resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        
            buffer = BytesIO()
            if img.mode in ('L', '1', 'LA', 'RGBA'):
                resized.save(buffer, format='PNG', optimize=True)
            else:
                if resized.mode not in ('RGB', 'CMYK'):
                    resized = resized.convert('RGB')
                resized.save(buffer, format='JPEG', quality=95, optimize=True)
        
            return buffer.getvalue()
        
        except Exception as e:
            logger.error(f"Ошибка масштабирования: {e}")
            return image_bytes


