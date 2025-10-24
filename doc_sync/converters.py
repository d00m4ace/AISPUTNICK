# doc_sync/converters.py
import re
import os
import sys
import logging
import tempfile
import subprocess
import shutil
from typing import Optional, Tuple
from io import BytesIO
from pathlib import Path

logger = logging.getLogger(__name__)

class HTMLToMarkdownConverter:
    """Конвертер HTML в Markdown"""
    
    @staticmethod
    def convert(html_content: str) -> str:
        """Конвертировать HTML в Markdown"""
        if not html_content:
            return ""
        
        # Используем markdownify для конвертации
        try:
            from markdownify import markdownify
            markdown = markdownify(
                html_content,
                heading_style="ATX",
                bullets="-",
                code_language="python",
                wrap=True,
                wrap_width=80
            )
        except ImportError:
            logger.warning("markdownify not installed, using simple conversion")
            markdown = HTMLToMarkdownConverter._simple_html_to_markdown(html_content)
        
        # Очистка результата
        markdown = HTMLToMarkdownConverter._clean_markdown(markdown)
        
        return markdown
    
    @staticmethod
    def _simple_html_to_markdown(html: str) -> str:
        """Простая конвертация HTML в Markdown без внешних библиотек"""
        # Удаляем скрипты и стили
        html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
        
        # Заголовки
        html = re.sub(r'<h1[^>]*>(.*?)</h1>', r'# \1', html, flags=re.IGNORECASE)
        html = re.sub(r'<h2[^>]*>(.*?)</h2>', r'## \1', html, flags=re.IGNORECASE)
        html = re.sub(r'<h3[^>]*>(.*?)</h3>', r'### \1', html, flags=re.IGNORECASE)
        html = re.sub(r'<h4[^>]*>(.*?)</h4>', r'#### \1', html, flags=re.IGNORECASE)
        html = re.sub(r'<h5[^>]*>(.*?)</h5>', r'##### \1', html, flags=re.IGNORECASE)
        html = re.sub(r'<h6[^>]*>(.*?)</h6>', r'###### \1', html, flags=re.IGNORECASE)
        
        # Форматирование
        html = re.sub(r'<b[^>]*>(.*?)</b>', r'**\1**', html, flags=re.IGNORECASE)
        html = re.sub(r'<strong[^>]*>(.*?)</strong>', r'**\1**', html, flags=re.IGNORECASE)
        html = re.sub(r'<i[^>]*>(.*?)</i>', r'*\1*', html, flags=re.IGNORECASE)
        html = re.sub(r'<em[^>]*>(.*?)</em>', r'*\1*', html, flags=re.IGNORECASE)
        html = re.sub(r'<code[^>]*>(.*?)</code>', r'`\1`', html, flags=re.IGNORECASE)
        
        # Списки
        html = re.sub(r'<li[^>]*>(.*?)</li>', r'- \1', html, flags=re.IGNORECASE)
        
        # Параграфы и переносы
        html = re.sub(r'<p[^>]*>(.*?)</p>', r'\1\n\n', html, flags=re.IGNORECASE | re.DOTALL)
        html = re.sub(r'<br[^>]*>', '\n', html, flags=re.IGNORECASE)
        html = re.sub(r'<hr[^>]*>', '\n---\n', html, flags=re.IGNORECASE)
        
        # Ссылки
        html = re.sub(r'<a[^>]*href="([^"]*)"[^>]*>(.*?)</a>', r'[\2](\1)', html, flags=re.IGNORECASE)
        
        # Удаляем все оставшиеся HTML теги
        html = re.sub(r'<[^>]+>', '', html)
        
        # Декодируем HTML entities
        html = html.replace('&nbsp;', ' ')
        html = html.replace('&lt;', '<')
        html = html.replace('&gt;', '>')
        html = html.replace('&amp;', '&')
        html = html.replace('&quot;', '"')
        
        return html
    
    @staticmethod
    def _clean_markdown(text: str) -> str:
        """Очистить и отформатировать markdown текст"""
        # Удаляем множественные пустые строки
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Удаляем пробелы в конце строк
        lines = text.split('\n')
        lines = [line.rstrip() for line in lines]
        text = '\n'.join(lines)
        
        # Удаляем HTML комментарии
        text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
        
        # Удаляем оставшиеся HTML теги
        text = re.sub(r'<[^>]+>', '', text)
        
        return text.strip()


class DocumentConverter:
    """Конвертер для документов Word/RTF/ODT"""
    
    SUPPORTED_EXTENSIONS = {'.docx', '.doc', '.rtf', '.odt'}
    
    @staticmethod
    def convert(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертировать документ в Markdown"""
        ext = os.path.splitext(filename.lower())[1]
        
        if ext not in DocumentConverter.SUPPORTED_EXTENSIONS:
            return False, ""
        
        # Сначала пробуем через pandoc
        if shutil.which("pandoc"):
            success, content = DocumentConverter._convert_via_pandoc(file_bytes, filename)
            if success:
                return True, content
        
        # Если pandoc не доступен, пробуем python-docx для .docx
        if ext == '.docx':
            success, content = DocumentConverter._convert_docx_native(file_bytes)
            if success:
                return True, content
        
        # Для RTF пробуем striprtf
        if ext == '.rtf':
            success, content = DocumentConverter._convert_rtf_native(file_bytes)
            if success:
                return True, content
        
        return False, ""
    
    @staticmethod
    def _convert_via_pandoc(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертировать через pandoc"""
        try:
            ext = os.path.splitext(filename.lower())[1]
            
            with tempfile.TemporaryDirectory() as td:
                in_path = os.path.join(td, f"input{ext}")
                out_path = os.path.join(td, "output.md")
                
                with open(in_path, "wb") as f:
                    f.write(file_bytes)
                
                # Используем расширенные параметры pandoc
                pandoc_args = [
                    "pandoc",
                    "-s",  # standalone
                    in_path,
                    "--wrap=auto",
                    "--columns=128",
                    "--reference-links",
                    "-t", "markdown",
                    "--extract-media", td,
                    "-o", out_path
                ]
                
                if ext == '.docx':
                    pandoc_args.insert(2, "--track-changes=accept")
                
                result = subprocess.run(
                    pandoc_args,
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                
                if result.returncode == 0:
                    with open(out_path, "r", encoding="utf-8", errors="replace") as f:
                        md = f.read()
                    
                    # Постобработка markdown
                    md = DocumentConverter._postprocess_markdown(md)
                    
                    if md.strip():
                        logger.info(f"Document converted via pandoc: {filename}")
                        return True, md
        
        except subprocess.TimeoutExpired:
            logger.error(f"Pandoc timeout for {filename}")
        except Exception as e:
            logger.error(f"Pandoc conversion error: {e}")
        
        return False, ""
    
    @staticmethod
    def _convert_docx_native(file_bytes: bytes) -> Tuple[bool, str]:
        """Конвертировать DOCX используя python-docx"""
        try:
            from docx import Document
            from io import BytesIO
            
            doc = Document(BytesIO(file_bytes))
            markdown = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Определяем стиль параграфа
                    if para.style.name.startswith('Heading 1'):
                        markdown.append(f"# {para.text}")
                    elif para.style.name.startswith('Heading 2'):
                        markdown.append(f"## {para.text}")
                    elif para.style.name.startswith('Heading 3'):
                        markdown.append(f"### {para.text}")
                    elif para.style.name.startswith('List'):
                        markdown.append(f"- {para.text}")
                    else:
                        # Обработка форматирования в тексте
                        text = para.text
                        for run in para.runs:
                            if run.bold:
                                text = text.replace(run.text, f"**{run.text}**")
                            elif run.italic:
                                text = text.replace(run.text, f"*{run.text}*")
                        markdown.append(text)
                    markdown.append("")
            
            # Обработка таблиц
            for table in doc.tables:
                markdown.append("")
                for i, row in enumerate(table.rows):
                    row_text = "| " + " | ".join([cell.text.strip() for cell in row.cells]) + " |"
                    markdown.append(row_text)
                    if i == 0:  # После заголовка добавляем разделитель
                        markdown.append("| " + " | ".join(["---"] * len(row.cells)) + " |")
                markdown.append("")
            
            content = "\n".join(markdown)
            if content.strip():
                logger.info("DOCX converted via python-docx")
                return True, content
                
        except Exception as e:
            logger.error(f"python-docx conversion error: {e}")
        
        return False, ""
    
    @staticmethod
    def _convert_rtf_native(file_bytes: bytes) -> Tuple[bool, str]:
        """Конвертировать RTF используя striprtf"""
        try:
            from striprtf.striprtf import rtf_to_text
            
            text = file_bytes.decode('utf-8', errors='replace')
            content = rtf_to_text(text)
            
            if content.strip():
                logger.info("RTF converted via striprtf")
                return True, content
                
        except ImportError:
            logger.warning("striprtf not installed")
        except Exception as e:
            logger.error(f"RTF conversion error: {e}")
        
        return False, ""
    
    @staticmethod
    def _postprocess_markdown(text: str) -> str:
        """Постобработка markdown после pandoc"""
        # Удаляем разметку {.mark} и подобную
        text = re.sub(r'\[([^\]]+)\]\{[^}]+\}', r'\1', text)
        text = re.sub(r'\{\.[\w\-]+\}', '', text)
        
        # Удаляем escape-символы
        text = re.sub(r'\\([#*_\[\]()>"|])', r'\1', text)
        
        # Удаляем двойные обратные слеши
        text = re.sub(r'\\\\', '', text)
        
        # Удаляем лишние пустые строки
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Удаляем пробелы в конце строк
        lines = text.split('\n')
        lines = [line.rstrip() for line in lines]
        text = '\n'.join(lines)
        
        # Удаляем мета-теги в начале документа если есть
        text = re.sub(r'^---\n.*?\n---\n', '', text, flags=re.DOTALL)
        
        return text.strip()


class SpreadsheetConverter:
    """Конвертер для Excel таблиц"""
    
    SUPPORTED_EXTENSIONS = {'.xls', '.xlsx'}
    
    @staticmethod
    def convert(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертировать таблицу в Markdown"""
        ext = os.path.splitext(filename.lower())[1]
        
        if ext not in SpreadsheetConverter.SUPPORTED_EXTENSIONS:
            return False, ""
        
        try:
            import pandas as pd
            from io import BytesIO
            
            # Читаем Excel файл
            xls = pd.ExcelFile(BytesIO(file_bytes))
            markdown = [f"# {filename}\n"]
            
            # Обрабатываем каждый лист
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name, dtype=object)
                df = df.fillna("")
                
                markdown.append(f"\n## {sheet_name}\n")
                markdown.append(f"*Строк: {len(df)}, Столбцов: {len(df.columns)}*\n\n")
                
                # Конвертируем в markdown таблицу
                rows = [list(df.columns)] + df.values.tolist()
                markdown.append(SpreadsheetConverter._rows_to_markdown(rows))
                markdown.append("\n")
            
            content = "\n".join(markdown)
            
            if content.strip():
                logger.info(f"Spreadsheet converted: {filename}")
                return True, content
                
        except ImportError:
            logger.error("pandas not installed for Excel conversion")
            # Пробуем альтернативный метод через openpyxl
            return SpreadsheetConverter._convert_via_openpyxl(file_bytes, filename)
        except Exception as e:
            logger.error(f"Excel conversion error: {e}")
        
        return False, ""
    
    @staticmethod
    def _convert_via_openpyxl(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Альтернативная конвертация через openpyxl"""
        try:
            from openpyxl import load_workbook
            from io import BytesIO
            
            wb = load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            markdown = [f"# {filename}\n"]
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                markdown.append(f"\n## {sheet_name}\n\n")
                
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else "" for cell in row])
                
                if rows:
                    markdown.append(SpreadsheetConverter._rows_to_markdown(rows))
                    markdown.append("\n")
            
            content = "\n".join(markdown)
            
            if content.strip():
                logger.info(f"Spreadsheet converted via openpyxl: {filename}")
                return True, content
                
        except Exception as e:
            logger.error(f"openpyxl conversion error: {e}")
        
        return False, ""
    
    @staticmethod
    def _rows_to_markdown(rows: list) -> str:
        """Преобразовать строки в markdown таблицу"""
        if not rows:
            return ""
        
        # Преобразуем все в строки
        str_rows = [[str(c) if c else "" for c in r] for r in rows]
        
        # Если таблица слишком большая, показываем только первые строки
        max_rows = 100
        if len(str_rows) > max_rows:
            str_rows = str_rows[:max_rows]
            truncated = True
        else:
            truncated = False
        
        # Вычисляем ширину колонок
        if str_rows:
            header = str_rows[0]
            body = str_rows[1:] if len(str_rows) > 1 else []
            
            widths = [len(h) for h in header]
            for r in body:
                for j, cell in enumerate(r):
                    if j < len(widths):
                        widths[j] = max(widths[j], len(cell[:50]))  # Ограничиваем ширину ячейки
            
            # Формируем таблицу
            def fmt_row(r):
                cells = []
                for i in range(len(widths)):
                    if i < len(r):
                        cell = r[i][:50]  # Ограничиваем длину содержимого
                        if len(r[i]) > 50:
                            cell += "..."
                        cells.append(cell.ljust(widths[i]))
                    else:
                        cells.append(" " * widths[i])
                return "| " + " | ".join(cells) + " |"
            
            parts = [fmt_row(header)]
            parts.append("| " + " | ".join(["-" * w for w in widths]) + " |")
            for r in body:
                parts.append(fmt_row(r))
            
            if truncated:
                parts.append("\n*...таблица обрезана (показаны первые 100 строк)*")
            
            return "\n".join(parts)
        
        return ""


class PresentationConverter:
    """Конвертер для презентаций PowerPoint"""
    
    SUPPORTED_EXTENSIONS = {'.pptx', '.ppt'}
    
    @staticmethod
    def convert(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертировать презентацию в Markdown"""
        ext = os.path.splitext(filename.lower())[1]
        
        if ext not in PresentationConverter.SUPPORTED_EXTENSIONS:
            return False, ""
        
        # Для .ppt используем только pandoc
        if ext == '.ppt':
            if shutil.which("pandoc"):
                success, content = PresentationConverter._convert_via_pandoc(file_bytes, filename)
                if success:
                    return True, content
            return False, ""
        
        # Для .pptx сначала пробуем python-pptx
        success, content = PresentationConverter._convert_via_python_pptx(file_bytes, filename)
        if success:
            return True, content
        
        # Если не получилось, пробуем pandoc
        if shutil.which("pandoc"):
            return PresentationConverter._convert_via_pandoc(file_bytes, filename)
        
        return False, ""
    
    @staticmethod
    def _convert_via_python_pptx(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертация через python-pptx"""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            prs = Presentation(BytesIO(file_bytes))
            markdown = [f"# {filename}\n"]
            
            for slide_num, slide in enumerate(prs.slides, 1):
                markdown.append(f"\n## Слайд {slide_num}\n")
                
                # Извлекаем текст из всех shape
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            # Определяем тип shape для форматирования
                            if shape.name.startswith("Title"):
                                markdown.append(f"### {text}\n")
                            elif shape.name.startswith("Subtitle"):
                                markdown.append(f"**{text}**\n")
                            else:
                                # Обрабатываем списки
                                lines = text.split('\n')
                                for line in lines:
                                    line = line.strip()
                                    if line:
                                        # Проверяем уровень отступа для вложенных списков
                                        if line.startswith('\t'):
                                            markdown.append(f"  - {line.strip()}")
                                        else:
                                            markdown.append(f"- {line}")
                                markdown.append("")
                    
                    # Обрабатываем таблицы
                    if shape.has_table:
                        markdown.append("\n")
                        table = shape.table
                        
                        # Заголовки таблицы
                        if table.rows:
                            headers = []
                            for cell in table.rows[0].cells:
                                headers.append(cell.text.strip())
                            
                            if headers:
                                markdown.append("| " + " | ".join(headers) + " |")
                                markdown.append("| " + " | ".join(["---"] * len(headers)) + " |")
                                
                                # Строки таблицы
                                for row in table.rows[1:]:
                                    cells = []
                                    for cell in row.cells:
                                        cells.append(cell.text.strip())
                                    markdown.append("| " + " | ".join(cells) + " |")
                        
                        markdown.append("\n")
                
                # Добавляем заметки к слайду если есть
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        markdown.append(f"\n**Заметки:**\n{notes}\n")
                
                # Разделитель между слайдами
                if slide_num < len(prs.slides):
                    markdown.append("\n---\n")
            
            content = "\n".join(markdown)
            
            if content.strip():
                logger.info(f"Presentation converted via python-pptx: {filename}")
                return True, content
                
        except ImportError:
            logger.warning("python-pptx not installed")
        except Exception as e:
            logger.error(f"python-pptx conversion error: {e}")
        
        return False, ""
    
    @staticmethod
    def _convert_via_pandoc(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертация через pandoc"""
        try:
            ext = os.path.splitext(filename.lower())[1]
            
            with tempfile.TemporaryDirectory() as td:
                in_path = os.path.join(td, f"input{ext}")
                out_path = os.path.join(td, "output.md")
                
                with open(in_path, "wb") as f:
                    f.write(file_bytes)
                
                subprocess.run(
                    ["pandoc", "-s", in_path, "--wrap=none", 
                     "-t", "markdown", "-o", out_path],
                    check=True, 
                    stdout=subprocess.PIPE, 
                    stderr=subprocess.PIPE,
                    timeout=30
                )
                
                with open(out_path, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                
                if content.strip():
                    logger.info(f"Presentation converted via pandoc: {filename}")
                    return True, content
                    
        except subprocess.TimeoutExpired:
            logger.error(f"Pandoc timeout for {filename}")
        except Exception as e:
            logger.error(f"Pandoc conversion error: {e}")
        
        return False, ""

class PDFConverter:
    """Конвертер PDF в Markdown"""
    
    @staticmethod
    def convert(file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """Конвертировать PDF в Markdown с улучшенной обработкой"""
        import tempfile
        import os
        
        # Создаем временный файл для работы с PDF
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
            tmp_path = tmp_file.name
            tmp_file.write(file_bytes)
            tmp_file.flush()
        
        try:
            logger.info(f"Processing PDF: {filename}")
            
            # Метод 1: MarkItDown
            success, content = PDFConverter._convert_via_markitdown_improved(tmp_path, filename)
            if success and content:
                return True, content
                        
            # Метод 2: OCR для сканов
            success, content = PDFConverter._convert_via_ocr(tmp_path, filename, file_bytes)
            if success and content:
                return True, content
            
            # Если все методы не сработали
            return PDFConverter._handle_conversion_failure(file_bytes)
            
        finally:
            # Всегда удаляем временный файл
            try:
                os.unlink(tmp_path)
                logger.debug(f"Removed temp file: {tmp_path}")
            except:
                pass
    
    @staticmethod
    def _convert_via_markitdown_improved(tmp_path: str, filename: str) -> Tuple[bool, str]:
        """Конвертация через MarkItDown с улучшенной обработкой"""
        try:
            from markitdown import MarkItDown
            
            markitdown = MarkItDown()
            result = markitdown.convert(tmp_path)
            
            if result and hasattr(result, 'text_content') and result.text_content:
                content = result.text_content
                # Удаляем строки "Content from the zip..." и "## File: ..."
                import re
                content = re.sub(r"(?m)^Content from the zip file.*\n?", "", content)
                content = re.sub(r"(?m)^## File:.*\n?", "", content)
                content = content.strip()
                if content and len(content) > 0:
                    logger.info(f"PDF extracted via MarkItDown: {filename}")
                    return True, content
        except ImportError:
            logger.debug("MarkItDown not available")
        except Exception as e:
            logger.warning(f"MarkItDown failed for {filename}: {e}")
        
        return False, ""
      
    @staticmethod
    def _convert_via_ocr(tmp_path: str, filename: str, file_bytes: bytes) -> Tuple[bool, str]:
        """OCR для сканированных PDF"""
        try:
            import pytesseract
            from PIL import Image
            import pdf2image
            
            logger.info(f"Attempting OCR for potential scanned PDF: {filename}")
            
            # Конвертируем PDF в изображения
            images = pdf2image.convert_from_path(tmp_path)
            
            text_parts = []
            
            for i, image in enumerate(images, 1):
                try:
                    # Применяем OCR
                    text = pytesseract.image_to_string(image, lang='rus+eng')
                    if text and text.strip():
                        text_parts.append(text.strip())
                        text_parts.append("\n\n")
                except Exception as e:
                    logger.warning(f"OCR failed for page {i}: {e}")
            
            if text_parts:
                content = "".join(text_parts).strip()
                logger.info(f"PDF extracted via OCR: {filename}")
                return True, content
                
        except ImportError:
            logger.debug("OCR tools not available (pytesseract/pdf2image)")
        except Exception as e:
            logger.warning(f"OCR failed for {filename}: {e}")
        
        return False, ""
    
    @staticmethod
    def _handle_conversion_failure(file_bytes: bytes) -> Tuple[bool, str]:
        """Обработка случая, когда все методы конвертации не сработали"""
        logger.error("All PDF extraction methods failed")
        
        error_msg = (
                f"*Could not extract text from this PDF file.*\n\n"
                f"The file may be corrupted, password-protected, or contain only images."
        )
        
        return False, error_msg

class UniversalFileConverter:
    """Универсальный конвертер файлов в Markdown"""
    
    def __init__(self):
        self.converters = {
            '.docx': DocumentConverter(),
            '.doc': DocumentConverter(),
            '.rtf': DocumentConverter(),
            '.odt': DocumentConverter(),
            '.xls': SpreadsheetConverter(),
            '.xlsx': SpreadsheetConverter(),
            '.pptx': PresentationConverter(),
            '.ppt': PresentationConverter(),
            '.pdf': PDFConverter(),
        }
        
        self.html_converter = HTMLToMarkdownConverter()
    
    def convert(self, file_bytes: bytes, filename: str) -> Tuple[bool, str]:
        """
        Универсальный метод конвертации файла в Markdown
        Returns: (success, markdown_content)
        """
        ext = os.path.splitext(filename.lower())[1]
        
        # Проверяем, есть ли конвертер для этого типа файла
        if ext in self.converters:
            converter = self.converters[ext]
            return converter.convert(file_bytes, filename)
        
        # Для HTML используем HTMLToMarkdownConverter
        if ext in ['.html', '.htm']:
            try:
                html_content = file_bytes.decode('utf-8', errors='replace')
                markdown = self.html_converter.convert(html_content)
                if markdown:
                    return True, markdown
            except Exception as e:
                logger.error(f"HTML conversion error: {e}")
        
        # Для текстовых файлов просто декодируем
        if ext in ['.txt', '.md', '.markdown', '.csv']:
            try:
                content = file_bytes.decode('utf-8', errors='replace')
                return True, content
            except Exception as e:
                logger.error(f"Text file decoding error: {e}")
        
        logger.warning(f"No converter available for file type: {ext}")
        return False, ""